"""
PPTX extraction using python-pptx
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path
from typing import Dict

def extract_from_pptx(pptx_path: str, out_dir: Path) -> Dict[str, any]:
    """Extract text and images from PPTX"""
    prs = Presentation(pptx_path)
    out_dir.mkdir(parents=True, exist_ok=True)
    full_text = []
    images = []

    for slide_index, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text"):
                    txt = shape.text.strip()
                    if txt:
                        slide_texts.append(txt)
            except Exception:
                pass

            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img = shape.image
                    ext = img.ext or "png"
                    blob = img.blob
                    filename = out_dir / f"slide{slide_index+1}_img_{len(images)+1}.{ext}"
                    with open(filename, "wb") as f:
                        f.write(blob)
                    images.append(filename)
            except Exception:
                pass

        if slide_texts:
            full_text.append(f"--- SLIDE {slide_index + 1} ---\n" + "\n".join(slide_texts) + "\n")

    return {"text": "\n".join(full_text), "images": images}